import uuid
import docx
from pptx import Presentation
from app.utils.text_cleaner import clean_text

class FileAgent:
    """Handles txt, docx, pptx extraction into unified content blocks."""

    def ingest(self, file_path: str, ext: str) -> dict:
        ext = ext.lower()
        file_id = f"file_{uuid.uuid4().hex[:8]}"
        text_blocks = []

        if ext == ".txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                raw = f.read()
            text_blocks.append({"source_id": file_id, "page": 1, "text": clean_text(raw)})

        elif ext == ".docx":
            doc = docx.Document(file_path)
            full = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])

            for table in doc.tables:
                for row in table.rows:
                    full += "\n" + " | ".join([c.text.strip() for c in row.cells])
            text_blocks.append({"source_id": file_id, "page": 1, "text": clean_text(full)})

        elif ext == ".pptx":
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        t = shape.text.strip()
                        if t:
                            texts.append(t)
                    if shape.has_table:
                        for row in shape.table.rows:
                            texts.append(" | ".join([c.text.strip() for c in row.cells]))
                if texts:
                    text_blocks.append({
                        "source_id": file_id,
                        "page": i + 1,
                        "text": clean_text("\n".join(texts))
                    })
        else:
            raise ValueError(f"Unsupported file type: {ext}")

        text_blocks = [b for b in text_blocks if b["text"]]

        return {
            "source_id": file_id,
            "type": "document",
            "content": text_blocks
        }
